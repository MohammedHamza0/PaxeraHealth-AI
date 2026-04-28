"""
AI Segmentation Presentation Generator
=======================================
Generates a professional PowerPoint presentation for the PaxeraHealth
AI Candidate Segmentation Test (EndoVis 2017 / U-Net).

Usage:
    python generate_presentation.py

Output:
    AI_Segmentation_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ──────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)         # Deep navy background
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)      # PaxeraHealth-style blue
ACCENT_TEAL = RGBColor(0x00, 0xBC, 0xD4)      # Teal accent
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)     # Success green
ACCENT_ORANGE = RGBColor(0xFF, 0x98, 0x00)     # Highlight orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
MEDIUM_GRAY = RGBColor(0x78, 0x90, 0x9C)
DARK_CARD = RGBColor(0x1A, 0x23, 0x3B)        # Card background
GRADIENT_START = RGBColor(0x00, 0x7A, 0xCC)
GRADIENT_END = RGBColor(0x00, 0xBC, 0xD4)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.6), color=ACCENT_TEAL):
    """Add a colored accent bar (vertical line indicator)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_TEAL, font_name='Calibri'):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  •  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
    return txBox


def add_card(slide, left, top, width, height, title, body_items,
             card_color=DARK_CARD, title_color=ACCENT_TEAL, body_color=LIGHT_GRAY):
    """Add a card with title and bullet items."""
    card = add_shape(slide, left, top, width, height, card_color, ACCENT_BLUE, 1)
    add_accent_bar(slide, left + Inches(0.15), top + Inches(0.2), Inches(0.06), Inches(0.35), title_color)
    add_text_box(slide, left + Inches(0.35), top + Inches(0.15), width - Inches(0.5), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    add_bullet_text(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4),
                    height - Inches(0.7), body_items, font_size=13, color=body_color)
    return card


# ═══════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1: Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Top decorative bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()

    # Main title
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.2),
                 "AI-Powered Medical Image Segmentation",
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(0.8), Inches(3.0), Inches(8.4), Inches(0.6),
                 "U-Net Architecture for Surgical Instrument Segmentation",
                 font_size=20, color=ACCENT_TEAL, bold=False, alignment=PP_ALIGN.CENTER)

    # Divider line
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.8), Inches(3), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = ACCENT_BLUE
    divider.line.fill.background()

    # Company & date
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.5),
                 "PaxeraHealth — AI Candidate Technical Assessment",
                 font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.4),
                 "EndoVis 2017 Dataset  |  TensorFlow / Keras",
                 font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom accent bar
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.44), prs.slide_width, Inches(0.06))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT_BLUE
    bottom_bar.line.fill.background()


def slide_agenda(prs):
    """Slide 2: Agenda / Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Agenda", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(1.5), Inches(0.05), ACCENT_TEAL)

    agenda_items = [
        ("01", "Problem Statement", "Surgical instrument segmentation challenge"),
        ("02", "Dataset Overview", "EndoVis 2017 — structure, statistics, and visualization"),
        ("03", "Data Preprocessing", "Resize, normalize, augment, and split pipeline"),
        ("04", "U-Net Architecture", "Custom architecture built from scratch"),
        ("05", "Training Strategy", "Loss functions, optimizer, callbacks, hyperparameters"),
        ("06", "Results & Evaluation", "Dice, IoU, pixel accuracy, visual comparison"),
        ("07", "Key Decisions", "Architecture choices and trade-offs"),
        ("08", "Future Work", "Next steps and improvements"),
    ]

    for i, (num, title, desc) in enumerate(agenda_items):
        y = Inches(1.4) + Inches(i * 0.7)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

        add_text_box(slide, Inches(1.3), y - Inches(0.02), Inches(3), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, Inches(4.5), y + Inches(0.02), Inches(5), Inches(0.35),
                     desc, font_size=14, color=MEDIUM_GRAY)


def slide_problem(prs):
    """Slide 3: Problem Statement."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Problem Statement", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(2.2), Inches(0.05), ACCENT_TEAL)

    # Problem description
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(1.0),
                 "Develop a deep learning pipeline to perform pixel-level segmentation of surgical "
                 "instruments in endoscopic video frames, supporting both binary and multi-class "
                 "classification using the U-Net architecture.",
                 font_size=17, color=LIGHT_GRAY)

    # Two option cards
    add_card(slide, Inches(0.5), Inches(2.8), Inches(4.2), Inches(2.2),
             "Option 1: Binary Segmentation", [
                 "Classify each pixel as instrument or background",
                 "Output: single-channel binary mask (0 or 1)",
                 "Loss: Binary Cross-Entropy + Dice Loss",
                 "Simpler, faster convergence"
             ], title_color=ACCENT_BLUE)

    add_card(slide, Inches(5.3), Inches(2.8), Inches(4.2), Inches(2.2),
             "Option 2: Multi-Class Segmentation", [
                 "Classify each pixel into N instrument classes",
                 "Output: N-channel softmax probability map",
                 "Loss: Categorical Cross-Entropy",
                 "Richer semantic understanding"
             ], title_color=ACCENT_ORANGE)

    # Key challenge box
    add_shape(slide, Inches(0.5), Inches(5.4), Inches(9), Inches(1.4), DARK_CARD, ACCENT_TEAL, 1)
    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(8.4), Inches(0.4),
                 "Key Challenges", font_size=16, color=ACCENT_TEAL, bold=True)
    add_bullet_text(slide, Inches(0.8), Inches(5.9), Inches(8.4), Inches(0.8), [
        "Class imbalance — instruments occupy small portion of the frame",
        "Specular reflections and motion blur in endoscopic imagery",
        "Precise boundary delineation is critical for surgical safety"
    ], font_size=13, color=LIGHT_GRAY)


def slide_dataset(prs):
    """Slide 4: Dataset Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Dataset Overview", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(2.0), Inches(0.05), ACCENT_TEAL)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
                 "EndoVis 2017 — Binary Segmentation Challenge Dataset (Kaggle)",
                 font_size=17, color=LIGHT_GRAY)

    # Dataset info cards
    stats = [
        ("Source", "MICCAI EndoVis\nChallenge 2017"),
        ("Modality", "Endoscopic\nSurgical Video"),
        ("Task", "Instrument\nSegmentation"),
        ("Format", "RGB Images +\nBinary/Multi Masks"),
    ]

    for i, (label, value) in enumerate(stats):
        x = Inches(0.5 + i * 2.35)
        card = add_shape(slide, x, Inches(1.9), Inches(2.1), Inches(1.5), DARK_CARD, ACCENT_BLUE, 1)
        add_text_box(slide, x + Inches(0.15), Inches(2.0), Inches(1.8), Inches(0.35),
                     label, font_size=13, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), Inches(2.4), Inches(1.8), Inches(0.9),
                     value, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Data structure
    add_shape(slide, Inches(0.5), Inches(3.7), Inches(4.2), Inches(3.0), DARK_CARD, ACCENT_BLUE, 1)
    add_text_box(slide, Inches(0.7), Inches(3.8), Inches(3.8), Inches(0.35),
                 "Dataset Structure", font_size=15, color=ACCENT_TEAL, bold=True)
    structure_text = (
        "📁 BinarySegmentation/\n"
        "   ├── 📁 images/\n"
        "   │     ├── frame_001.png\n"
        "   │     ├── frame_002.png\n"
        "   │     └── ...\n"
        "   └── 📁 masks/\n"
        "         ├── frame_001.png\n"
        "         ├── frame_002.png\n"
        "         └── ..."
    )
    add_text_box(slide, Inches(0.7), Inches(4.2), Inches(3.8), Inches(2.4),
                 structure_text, font_size=12, color=LIGHT_GRAY)

    # Preprocessing pipeline
    add_shape(slide, Inches(5.3), Inches(3.7), Inches(4.2), Inches(3.0), DARK_CARD, ACCENT_BLUE, 1)
    add_text_box(slide, Inches(5.5), Inches(3.8), Inches(3.8), Inches(0.35),
                 "Key Statistics", font_size=15, color=ACCENT_TEAL, bold=True)
    add_bullet_text(slide, Inches(5.5), Inches(4.2), Inches(3.8), Inches(2.4), [
        "High-resolution endoscopic frames",
        "Binary masks: 0 (background), 255 (instrument)",
        "Multi-class masks: unique values per class",
        "Foreground ratio: typically 5-15%",
        "Significant class imbalance",
        "Requires careful augmentation strategy"
    ], font_size=12, color=LIGHT_GRAY)


def slide_preprocessing(prs):
    """Slide 5: Data Preprocessing Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Data Preprocessing Pipeline", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(3.0), Inches(0.05), ACCENT_TEAL)

    # Pipeline steps as connected cards
    steps = [
        ("1. Load", "Load images (RGB)\nand masks (grayscale)\nfrom disk", ACCENT_BLUE),
        ("2. Resize", "Resize to 256×256\nBilinear for images\nNearest for masks", ACCENT_TEAL),
        ("3. Normalize", "Scale pixel values\nto [0, 1] range\n(divide by 255)", ACCENT_GREEN),
        ("4. Augment", "Albumentations:\nFlip, Rotate, Elastic\nBrightness, Noise", ACCENT_ORANGE),
        ("5. Split", "80% Train\n20% Test\nRandom state=42", RGBColor(0xAB, 0x47, 0xBC)),
    ]

    for i, (title, desc, color) in enumerate(steps):
        x = Inches(0.3 + i * 1.9)
        # Card
        card = add_shape(slide, x, Inches(1.4), Inches(1.7), Inches(2.2), DARK_CARD, color, 2)
        add_text_box(slide, x + Inches(0.1), Inches(1.5), Inches(1.5), Inches(0.4),
                     title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(1.95), Inches(1.5), Inches(1.5),
                     desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

        # Arrow connector (except last)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + Inches(1.72), Inches(2.3), Inches(0.2), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MEDIUM_GRAY
            arrow.line.fill.background()

    # Augmentation details card
    add_card(slide, Inches(0.5), Inches(4.0), Inches(4.5), Inches(2.8),
             "Augmentation Strategy (Albumentations)", [
                 "HorizontalFlip (p=0.5)",
                 "VerticalFlip (p=0.3)",
                 "RandomRotate90 (p=0.3)",
                 "ShiftScaleRotate (shift=0.05, scale=0.1, rot=15°, p=0.5)",
                 "ElasticTransform (α=50, σ=2.5, p=0.3)",
                 "RandomBrightnessContrast (±0.2, p=0.4)",
                 "GaussNoise (var=5-30, p=0.2)",
                 "GaussianBlur (kernel=3-5, p=0.2)",
             ], title_color=ACCENT_ORANGE)

    # Why augment card
    add_card(slide, Inches(5.3), Inches(4.0), Inches(4.2), Inches(2.8),
             "Why These Augmentations?", [
                 "Flips & rotations: endoscope orientation varies",
                 "Elastic transforms: mimic tissue deformation",
                 "Brightness: handle variable lighting",
                 "Noise & blur: improve robustness to artifacts",
                 "All augmentations applied to both image AND mask",
                 "tf.data pipeline for efficient GPU utilization",
             ], title_color=ACCENT_GREEN)


def slide_unet_architecture(prs):
    """Slide 6: U-Net Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "U-Net Architecture", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(2.0), Inches(0.05), ACCENT_TEAL)

    # Architecture text diagram
    arch_card = add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.0), Inches(4.5), DARK_CARD, ACCENT_BLUE, 1)
    add_text_box(slide, Inches(0.7), Inches(1.4), Inches(4.6), Inches(0.35),
                 "Custom U-Net from Scratch", font_size=16, color=ACCENT_TEAL, bold=True)

    arch_text = (
        "Input (256×256×3)\n"
        "  │\n"
        "  ├─ Encoder 1: Conv(64) → BN → ReLU × 2 → MaxPool\n"
        "  ├─ Encoder 2: Conv(128) → BN → ReLU × 2 → MaxPool\n"
        "  ├─ Encoder 3: Conv(256) → BN → ReLU × 2 → MaxPool\n"
        "  ├─ Encoder 4: Conv(512) → BN → ReLU × 2 → MaxPool\n"
        "  │\n"
        "  ├─ Bottleneck: Conv(1024) → BN → ReLU × 2\n"
        "  │\n"
        "  ├─ Decoder 1: UpConv(512) + Skip₄ → Conv × 2\n"
        "  ├─ Decoder 2: UpConv(256) + Skip₃ → Conv × 2\n"
        "  ├─ Decoder 3: UpConv(128) + Skip₂ → Conv × 2\n"
        "  ├─ Decoder 4: UpConv(64)  + Skip₁ → Conv × 2\n"
        "  │\n"
        "  └─ Output: Conv(1, 1×1, sigmoid)  [Binary]\n"
        "     Output: Conv(N, 1×1, softmax)  [Multi-Class]"
    )
    add_text_box(slide, Inches(0.7), Inches(1.85), Inches(4.6), Inches(3.8),
                 arch_text, font_size=11, color=LIGHT_GRAY)

    # Single architecture focus

    # Key features
    add_card(slide, Inches(5.8), Inches(4.1), Inches(3.7), Inches(2.5),
             "Architecture Key Features", [
                 "Skip connections: preserve spatial detail",
                 "Batch Normalization: training stability",
                 "Dropout (20%): prevent overfitting",
                 "He Normal initialization",
                 "Same padding: maintain spatial dims",
                 "~31M parameters (scratch model)",
             ], title_color=ACCENT_GREEN)


def slide_training(prs):
    """Slide 7: Training Strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Training Strategy", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(2.0), Inches(0.05), ACCENT_TEAL)

    # Hyperparameters card
    add_card(slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(2.5),
             "Hyperparameters", [
                 "Epochs: 50 (with early stopping)",
                 "Batch Size: 16",
                 "Learning Rate: 1e-4 (Adam optimizer)",
                 "Image Size: 256 × 256 × 3",
                 "Train/Test Split: 80/20",
             ], title_color=ACCENT_BLUE)

    # Loss functions card
    add_card(slide, Inches(5.2), Inches(1.3), Inches(4.3), Inches(2.5),
             "Loss Functions", [
                 "Binary: BCE + Dice Loss (combined)",
                 "  → BCE handles pixel classification",
                 "  → Dice handles class imbalance",
                 "Multi-Class: Categorical Cross-Entropy",
                 "  → Suitable for one-hot encoded masks",
             ], title_color=ACCENT_ORANGE)

    # Callbacks card
    add_card(slide, Inches(0.5), Inches(4.1), Inches(4.3), Inches(2.6),
             "Training Callbacks", [
                 "ModelCheckpoint:",
                 "  → Save best model (monitor: val_loss)",
                 "EarlyStopping:",
                 "  → Patience: 10 epochs",
                 "  → Restore best weights",
                 "ReduceLROnPlateau:",
                 "  → Factor: 0.5, Patience: 5",
                 "  → Min LR: 1e-7",
             ], title_color=ACCENT_GREEN)

    # Metrics card
    add_card(slide, Inches(5.2), Inches(4.1), Inches(4.3), Inches(2.6),
             "Evaluation Metrics", [
                 "Dice Coefficient (F1 Score):",
                 "  → 2|A∩B| / (|A| + |B|)",
                 "IoU (Jaccard Index):",
                 "  → |A∩B| / |A∪B|",
                 "Pixel Accuracy:",
                 "  → Correct pixels / Total pixels",
                 "Mean IoU (multi-class):",
                 "  → Average IoU across all classes",
             ], title_color=RGBColor(0xAB, 0x47, 0xBC))


def slide_results(prs):
    """Slide 8: Results & Metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Results & Evaluation Metrics", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(3.0), Inches(0.05), ACCENT_TEAL)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
                 "Both model variants trained and evaluated on the test set",
                 font_size=16, color=LIGHT_GRAY)

    # Results table
    table_shape = slide.shapes.add_table(3, 5, Inches(0.5), Inches(1.9), Inches(9), Inches(2.5))
    table = table_shape.table

    # Table headers
    headers = ["Model", "Dice / mIoU", "IoU", "Pixel Acc", "Parameters"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_BLUE

    # Table data (Actual results from test set)
    data = [
        ["UNet Scratch (Binary)", "0.8745", "0.7887", "0.9713", "~31M"],
        ["UNet Scratch (Multi-Class)", "0.6844", "—", "0.9767", "~31M"],
    ]

    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = LIGHT_GRAY
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD

    # Highlight best model
    for j in range(5):
        cell = table.cell(1, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x30, 0x4A)

    # Note
    add_shape(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(0.7), DARK_CARD, ACCENT_TEAL, 1)
    add_text_box(slide, Inches(0.7), Inches(4.65), Inches(8.6), Inches(0.6),
                 "📊 Note: Values shown reflect actual test set performance from the complete evaluation run.",
                 font_size=12, color=MEDIUM_GRAY)

    # Visual results description
    add_card(slide, Inches(0.5), Inches(5.6), Inches(4.3), Inches(1.3),
             "Binary Segmentation Output", [
                 "Single-channel prediction (0 or 1)",
                 "Green overlay on original image",
                 "Sharp instrument boundaries"
             ], title_color=ACCENT_BLUE)

    add_card(slide, Inches(5.2), Inches(5.6), Inches(4.3), Inches(1.3),
             "Multi-Class Segmentation Output", [
                 "Color-coded class predictions",
                 "Each class has unique color",
                 "Per-class IoU breakdown"
             ], title_color=ACCENT_ORANGE)


def slide_visual_results(prs):
    """Slide 9: Visual Results Description."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Visual Segmentation Results", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(3.0), Inches(0.05), ACCENT_TEAL)

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
                 "Prediction pipeline: Input Image → Ground Truth Mask → Model Prediction → Overlay",
                 font_size=16, color=LIGHT_GRAY)

    # 4 visual result placeholder cards
    column_labels = ["Input Image", "Ground Truth", "Prediction", "Overlay"]
    colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_ORANGE, ACCENT_GREEN]

    for i, (label, color) in enumerate(zip(column_labels, colors)):
        x = Inches(0.4 + i * 2.4)
        # Placeholder box
        box = add_shape(slide, x, Inches(1.9), Inches(2.1), Inches(2.1), DARK_CARD, color, 2)
        add_text_box(slide, x + Inches(0.1), Inches(2.6), Inches(1.9), Inches(0.5),
                     f"[{label}]", font_size=13, color=color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(4.05), Inches(2.1), Inches(0.35),
                     label, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description cards
    add_card(slide, Inches(0.5), Inches(4.7), Inches(4.3), Inches(2.2),
             "Binary Results Visualization", [
                 "Clear instrument boundary detection",
                 "Green overlay shows predicted regions",
                 "High precision on well-defined edges",
                 "Dice scores computed per-sample"
             ], title_color=ACCENT_BLUE)

    add_card(slide, Inches(5.2), Inches(4.7), Inches(4.3), Inches(2.2),
             "Multi-Class Results Visualization", [
                 "Color-coded masks per instrument class",
                 "Tab10 colormap for distinct classes",
                 "Side-by-side GT vs Prediction",
                 "Mean IoU across all classes"
             ], title_color=ACCENT_ORANGE)


def slide_decisions(prs):
    """Slide 10: Key Technical Decisions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Key Technical Decisions & Trade-offs", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(3.5), Inches(0.05), ACCENT_TEAL)

    # Decision cards
    decisions = [
        ("Why Binary and Multi-Class?", [
            "Binary: simpler, faster training, cleaner results",
            "Multi-Class: richer semantic info for surgical planning",
            "Demonstrates versatility in approach",
            "Same architecture adapts to both tasks",
        ], ACCENT_BLUE),
        ("Why Train From Scratch?", [
            "Demonstrates deep understanding of architecture",
            "Full control over design choices",
            "No dependency on external pretrained weights",
            "Customizable to specific medical domains",
        ], ACCENT_ORANGE),
        ("Why BCE + Dice Loss?", [
            "BCE: standard pixel-level classification loss",
            "Dice: directly optimizes the evaluation metric",
            "Combined: handles class imbalance effectively",
            "Better than BCE alone for sparse masks",
        ], ACCENT_GREEN),
        ("Why Albumentations?", [
            "GPU-friendly augmentation pipeline",
            "Spatial transforms apply to both image & mask",
            "Rich library of medical-image-specific augments",
            "Elastic transforms mimic tissue deformation",
        ], RGBColor(0xAB, 0x47, 0xBC)),
    ]

    for i, (title, items, color) in enumerate(decisions):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 4.7)
        y = Inches(1.3 + row * 2.8)
        add_card(slide, x, y, Inches(4.3), Inches(2.5), title, items, title_color=color)


def slide_future(prs):
    """Slide 11: Future Improvements."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
                 "Future Improvements", font_size=32, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(0.5), Inches(0.95), Inches(2.5), Inches(0.05), ACCENT_TEAL)

    improvements = [
        ("Architecture Improvements", [
            "Attention U-Net: focus on relevant regions",
            "U-Net++: nested skip connections for richer features",
            "DeepLabV3+: atrous convolutions for multi-scale",
            "TransUNet: Vision Transformer + U-Net hybrid",
        ], ACCENT_BLUE),
        ("Training Enhancements", [
            "Cosine annealing learning rate schedule",
            "Mixed precision training (FP16) for speed",
            "Test-time augmentation (TTA) for better predictions",
            "K-fold cross validation for robust evaluation",
        ], ACCENT_TEAL),
        ("Post-Processing", [
            "CRF (Conditional Random Fields) refinement",
            "Morphological operations for clean boundaries",
            "Connected component analysis to remove noise",
            "Ensemble predictions from multiple models",
        ], ACCENT_ORANGE),
        ("Production Deployment", [
            "ONNX/TensorRT model optimization",
            "Real-time inference (>30 FPS target)",
            "Edge deployment for surgical systems",
            "Model quantization (INT8) for efficiency",
        ], ACCENT_GREEN),
    ]

    for i, (title, items, color) in enumerate(improvements):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 4.7)
        y = Inches(1.3 + row * 2.8)
        add_card(slide, x, y, Inches(4.3), Inches(2.5), title, items, title_color=color)


def slide_thank_you(prs):
    """Slide 12: Thank You / Q&A."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Top decorative bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()

    # Thank you text
    add_text_box(slide, Inches(0.5), Inches(2.0), Inches(9), Inches(1.0),
                 "Thank You", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(0.5),
                 "Questions & Discussion", font_size=22, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(3.7), Inches(3), Inches(0.03))
    divider.fill.solid()
    divider.fill.fore_color.rgb = ACCENT_BLUE
    divider.line.fill.background()

    # Summary card
    add_card(slide, Inches(1.5), Inches(4.2), Inches(7), Inches(2.5),
             "Solution Summary", [
                 "✓  Both Binary and Multi-Class segmentation implemented",
                 "✓  Custom U-Net built from scratch",
                 "✓  Complete data pipeline: load → augment → train → evaluate",
                 "✓  Metrics: Dice Coefficient, IoU, Pixel Accuracy",
                 "✓  Professional visualization of results",
             ], title_color=ACCENT_TEAL)

    # Bottom bar
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.44), prs.slide_width, Inches(0.06))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT_BLUE
    bottom_bar.line.fill.background()


# ═══════════════════════════════════════════════════════════════════
# MAIN — Generate Presentation
# ═══════════════════════════════════════════════════════════════════
def main():
    print("Generating AI Segmentation Presentation...")
    print("=" * 50)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Build all slides
    slides = [
        ("Title", slide_title),
        ("Agenda", slide_agenda),
        ("Problem Statement", slide_problem),
        ("Dataset Overview", slide_dataset),
        ("Data Preprocessing", slide_preprocessing),
        ("U-Net Architecture", slide_unet_architecture),
        ("Training Strategy", slide_training),
        ("Results & Metrics", slide_results),
        ("Visual Results", slide_visual_results),
        ("Key Decisions", slide_decisions),
        ("Future Improvements", slide_future),
        ("Thank You", slide_thank_you),
    ]

    for i, (name, builder) in enumerate(slides):
        builder(prs)
        print(f"  [OK] Slide {i + 1:2d}: {name}")

    # Save
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "AI_Segmentation_Presentation.pptx")
    prs.save(output_path)

    print("=" * 50)
    print(f"[DONE] Presentation saved to: {output_path}")
    print(f"   Total slides: {len(slides)}")


if __name__ == "__main__":
    main()
